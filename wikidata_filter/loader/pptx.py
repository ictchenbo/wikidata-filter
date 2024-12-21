import os
from typing import Iterable, Any
from wikidata_filter.loader.file import BinaryFile


try:
    from pptx import Presentation
except:
    print("Error! you need to install pptx first: pip install pptx")
    raise ImportError("pptx not installed")


class PPTX(BinaryFile):
    """基于python-pptx读取pptx文件"""
    def __init__(self, input_file: str, max_pages: int = 0, **kwargs):
        super().__init__(input_file, auto_open=False, **kwargs)
        self.max_pages = max_pages

    def iter(self) -> Iterable[Any]:
        pts1 = Presentation(self.filename)

        for pageno, slide in enumerate(pts1.slides):
            if 0 < self.max_pages <= pageno:
                break
            paragraphs = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        part = []
                        for run in paragraph.runs:
                            paragraphs.append(run.text)
            yield {
                'format': 'pptx',
                'page': pageno,
                'content': paragraphs
            }


class PPT(BinaryFile):
    """基于libreoffice读取ppt文件，转换成pptx文件"""
    def __init__(self, input_file, max_pages=0, **kwargs):
        super().__init__(input_file, max_pages=max_pages, auto_open=False, **kwargs)
        pptx = input_file + "x"
        if not os.path.exists(pptx):
            out_path = os.path.dirname(self.filename)
            try:
                self.ppt2pptx(self.filename, out_path)
            except:
                print('ppt2pptx failed! make sure you have install libreoffice7.6')
                raise FileNotFoundError("libreoffice7.6 not found")

        self.filename = pptx
        self.instream = open(self.filename, 'rb')

    def ppt2pptx(self, doc: str, output_dir: str):
        os.system(f'libreoffice7.6 --headless --convert-to pptx "{doc}" --outdir "{output_dir}"')
